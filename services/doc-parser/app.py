"""Small, local-first document-to-Markdown service for Neo Chat.

The service intentionally does not perform OCR or call a third-party API. It
accepts one file, parses it in a temporary directory, and removes the source
before returning. The supported formats cover the common office/document
attachments without bringing a full browser or LibreOffice stack into the
image.
"""

from __future__ import annotations

import os
import tempfile
import zipfile
from pathlib import Path
from typing import Any, Iterable

import fitz  # PyMuPDF
import xlrd
from docx import Document
from fastapi import FastAPI, File, HTTPException, UploadFile
from openpyxl import load_workbook
from pptx import Presentation


def _env_int(name: str, default: int, minimum: int, maximum: int) -> int:
    try:
        value = int(os.getenv(name, str(default)))
    except (TypeError, ValueError):
        return default
    return min(max(value, minimum), maximum)


MAX_FILE_BYTES = _env_int(
    "DOC_PARSER_MAX_FILE_BYTES", 50 * 1024 * 1024, 1, 200 * 1024 * 1024
)
MAX_MARKDOWN_CHARS = _env_int(
    "DOC_PARSER_MAX_MARKDOWN_CHARS", 2_000_000, 1_000, 20_000_000
)
MAX_PDF_PAGES = _env_int("DOC_PARSER_MAX_PDF_PAGES", 500, 1, 5_000)
MAX_ARCHIVE_ENTRIES = _env_int("DOC_PARSER_MAX_ARCHIVE_ENTRIES", 2_000, 1, 20_000)
MAX_ARCHIVE_BYTES = _env_int(
    "DOC_PARSER_MAX_ARCHIVE_BYTES", 100 * 1024 * 1024, 1, 500 * 1024 * 1024
)
MAX_COMPRESSION_RATIO = _env_int("DOC_PARSER_MAX_COMPRESSION_RATIO", 100, 1, 1_000)
MAX_TABLE_ROWS = _env_int("DOC_PARSER_MAX_TABLE_ROWS", 50_000, 1, 500_000)
MAX_TABLE_COLUMNS = _env_int("DOC_PARSER_MAX_TABLE_COLUMNS", 200, 1, 1_000)

TEXT_EXTENSIONS = {
    ".c",
    ".cpp",
    ".css",
    ".csv",
    ".go",
    ".graphql",
    ".html",
    ".htm",
    ".java",
    ".js",
    ".json",
    ".jsx",
    ".log",
    ".md",
    ".php",
    ".py",
    ".rb",
    ".rs",
    ".sh",
    ".sql",
    ".svg",
    ".tex",
    ".ts",
    ".tsx",
    ".txt",
    ".xml",
    ".yaml",
    ".yml",
}


class ParserError(Exception):
    def __init__(self, message: str, status_code: int = 422):
        super().__init__(message)
        self.status_code = status_code


def _safe_filename(filename: str | None) -> str:
    name = Path(filename or "document").name.strip()
    if not name or name in {".", ".."}:
        return "document"
    return name[:255]


async def _save_upload(upload: UploadFile, suffix: str) -> Path:
    path: Path | None = None
    total = 0
    try:
        with tempfile.NamedTemporaryFile(
            prefix="neo-doc-", suffix=suffix, delete=False
        ) as target:
            path = Path(target.name)
            while True:
                chunk = await upload.read(1024 * 1024)
                if not chunk:
                    break
                total += len(chunk)
                if total > MAX_FILE_BYTES:
                    raise ParserError(
                        f"Document exceeds the {MAX_FILE_BYTES} byte limit", 413
                    )
                target.write(chunk)
    except Exception:
        if path:
            path.unlink(missing_ok=True)
        raise
    finally:
        await upload.close()

    if total == 0:
        path.unlink(missing_ok=True)
        raise ParserError("Document file is empty", 400)
    return path


def _validate_zip(path: Path) -> None:
    try:
        with zipfile.ZipFile(path) as archive:
            entries = archive.infolist()
            if len(entries) > MAX_ARCHIVE_ENTRIES:
                raise ParserError("Document archive contains too many entries", 413)

            total_size = 0
            for info in entries:
                name = info.filename.replace("\\", "/")
                if name.startswith("/") or ".." in Path(name).parts:
                    raise ParserError("Document archive contains an unsafe path", 422)
                total_size += info.file_size
                if total_size > MAX_ARCHIVE_BYTES:
                    raise ParserError("Document archive expands to too much data", 413)
                compressed_size = max(info.compress_size, 1)
                if info.file_size > compressed_size * MAX_COMPRESSION_RATIO:
                    raise ParserError("Document archive compression ratio is too high", 422)
    except ParserError:
        raise
    except (OSError, zipfile.BadZipFile) as error:
        raise ParserError("The document package is not a valid ZIP archive", 422) from error


def _clip(value: Any, limit: int = 8_000) -> str:
    if value is None:
        return ""
    text = str(value).replace("\x00", "").strip()
    return text[:limit]


def _escape_cell(value: Any) -> str:
    return _clip(value).replace("|", "\\|").replace("\r", " ").replace("\n", "<br>")


def _markdown_table(rows: Iterable[Iterable[Any]]) -> str:
    normalized: list[list[str]] = []
    for raw_row in rows:
        row = [_escape_cell(value) for value in list(raw_row)[:MAX_TABLE_COLUMNS]]
        while row and not row[-1]:
            row.pop()
        if row:
            normalized.append(row)
        if len(normalized) >= MAX_TABLE_ROWS:
            break

    if not normalized:
        return ""
    width = max(len(row) for row in normalized)
    normalized = [row + [""] * (width - len(row)) for row in normalized]
    header = normalized[0]
    lines = [
        "| " + " | ".join(header) + " |",
        "| " + " | ".join("---" for _ in range(width)) + " |",
    ]
    lines.extend("| " + " | ".join(row) + " |" for row in normalized[1:])
    return "\n".join(lines)


def _decode_text(path: Path) -> str:
    raw = path.read_bytes()
    for encoding in ("utf-8-sig", "utf-8", "gb18030", "cp1252", "latin-1"):
        try:
            return raw.decode(encoding)
        except UnicodeDecodeError:
            continue
    return raw.decode("utf-8", errors="replace")


def _parse_text(path: Path) -> str:
    return _decode_text(path)


def _parse_pdf(path: Path) -> str:
    try:
        document = fitz.open(str(path))
    except Exception as error:
        raise ParserError("Unable to open the PDF document", 422) from error

    if document.page_count > MAX_PDF_PAGES:
        document.close()
        raise ParserError("PDF contains too many pages", 413)

    pages: list[str] = []
    try:
        for index, page in enumerate(document, start=1):
            text = page.get_text("text").strip()
            if text:
                pages.append(f"## Page {index}\n\n{text}")
    finally:
        document.close()

    if not pages:
        raise ParserError(
            "No text layer was found in this PDF. Scanned PDFs require an OCR-enabled parser.",
            422,
        )
    return "\n\n".join(pages)


def _parse_docx(path: Path) -> str:
    _validate_zip(path)
    try:
        document = Document(str(path))
    except Exception as error:
        raise ParserError("Unable to open the DOCX document", 422) from error

    blocks: list[str] = []
    for paragraph in document.paragraphs:
        text = paragraph.text.strip()
        if not text:
            continue
        style_name = (paragraph.style.name or "").lower()
        if "heading" in style_name:
            level = next(
                (int(char) for char in reversed(style_name) if char.isdigit()), 2
            )
            blocks.append(f"{'#' * min(max(level, 1), 6)} {text}")
        else:
            blocks.append(text)

    for table in document.tables:
        table_markdown = _markdown_table(
            ([cell.text for cell in row.cells] for row in table.rows)
        )
        if table_markdown:
            blocks.append(table_markdown)
    return "\n\n".join(blocks)


def _parse_xlsx(path: Path) -> str:
    _validate_zip(path)
    try:
        workbook = load_workbook(str(path), read_only=True, data_only=True)
    except Exception as error:
        raise ParserError("Unable to open the XLSX workbook", 422) from error

    blocks: list[str] = []
    try:
        for sheet in workbook.worksheets:
            table = _markdown_table(sheet.iter_rows(values_only=True))
            if table:
                blocks.append(f"## {sheet.title}\n\n{table}")
    finally:
        workbook.close()
    return "\n\n".join(blocks)


def _parse_xls(path: Path) -> str:
    try:
        workbook = xlrd.open_workbook(str(path), on_demand=True)
    except Exception as error:
        raise ParserError("Unable to open the XLS workbook", 422) from error

    blocks: list[str] = []
    try:
        for sheet in workbook.sheets():
            rows = (
                sheet.row_values(index)
                for index in range(min(sheet.nrows, MAX_TABLE_ROWS))
            )
            table = _markdown_table(rows)
            if table:
                blocks.append(f"## {sheet.name}\n\n{table}")
    finally:
        workbook.release_resources()
    return "\n\n".join(blocks)


def _parse_pptx(path: Path) -> str:
    _validate_zip(path)
    try:
        presentation = Presentation(str(path))
    except Exception as error:
        raise ParserError("Unable to open the PPTX presentation", 422) from error

    slides: list[str] = []
    for slide_number, slide in enumerate(presentation.slides, start=1):
        blocks = [f"## Slide {slide_number}"]
        for shape in slide.shapes:
            if getattr(shape, "has_table", False):
                table = _markdown_table(
                    ([cell.text for cell in row.cells] for row in shape.table.rows)
                )
                if table:
                    blocks.append(table)
            else:
                text = getattr(shape, "text", "").strip()
                if text:
                    blocks.append(text)
        if len(blocks) > 1:
            slides.append("\n\n".join(blocks))
    return "\n\n".join(slides)


def _parse_file(path: Path, original_name: str, content_type: str | None) -> str:
    extension = Path(original_name).suffix.lower()
    mime_type = (content_type or "").split(";", 1)[0].strip().lower()
    if not extension:
        extension = {
            "application/pdf": ".pdf",
            "application/msword": ".doc",
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document": ".docx",
            "application/vnd.ms-excel": ".xls",
            "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": ".xlsx",
            "application/vnd.ms-powerpoint": ".ppt",
            "application/vnd.openxmlformats-officedocument.presentationml.presentation": ".pptx",
        }.get(mime_type, "")
    if extension in {".png", ".jpg", ".jpeg", ".gif", ".webp", ".bmp", ".tiff"}:
        raise ParserError(
            "Image OCR is not enabled in the lightweight local parser",
            415,
        )
    if extension in {".doc", ".ppt", ".xlsb"}:
        raise ParserError(
            f"{extension} files are not supported; convert the file to DOCX, PPTX, or XLSX first",
            415,
        )
    if extension in TEXT_EXTENSIONS or mime_type.startswith("text/"):
        return _parse_text(path)
    if extension == ".pdf":
        return _parse_pdf(path)
    if extension == ".docx":
        return _parse_docx(path)
    if extension == ".xlsx":
        return _parse_xlsx(path)
    if extension == ".xls":
        return _parse_xls(path)
    if extension == ".pptx":
        return _parse_pptx(path)
    raise ParserError(
        "Unsupported document type. Supported formats: PDF, DOCX, XLSX, XLS, PPTX, and text files.",
        415,
    )


def _finalize_markdown(markdown: str) -> str:
    result = markdown.replace("\x00", "").strip()
    if len(result) > MAX_MARKDOWN_CHARS:
        raise ParserError("Parsed Markdown is too large", 413)
    if not result:
        raise ParserError("No text content was extracted from the document", 422)
    return result


app = FastAPI(title="Neo Chat local document parser", docs_url=None, redoc_url=None)


@app.get("/healthz")
async def healthz() -> dict[str, str]:
    return {"status": "ok", "parser": "neo-chat-doc-parser", "version": "1"}


@app.post("/parse")
async def parse(file: UploadFile = File(...)) -> dict[str, str]:
    filename = _safe_filename(file.filename)
    suffix = Path(filename).suffix.lower()[:16]
    path: Path | None = None
    try:
        path = await _save_upload(file, suffix)
        markdown = _parse_file(path, filename, file.content_type)
        return {"markdown": _finalize_markdown(markdown)}
    except ParserError as error:
        raise HTTPException(status_code=error.status_code, detail=str(error)) from error
    except HTTPException:
        raise
    except Exception as error:
        # Do not return parser/library internals to the caller. The container
        # log still contains a useful exception for an operator to diagnose.
        import logging

        logging.getLogger(__name__).exception("document parse failed")
        raise HTTPException(status_code=422, detail="Unable to parse the document") from error
    finally:
        if path:
            path.unlink(missing_ok=True)
